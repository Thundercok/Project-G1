import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Set to 16:9 Widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette (Project G1 Retro Sci-Fi Theme)
    BG_DARK = RGBColor(18, 22, 28)       # Deep slate blue/black
    BG_CARD = RGBColor(28, 34, 44)       # Lighter card container
    ACCENT_AMBER = RGBColor(255, 176, 0) # GoldSrc Amber
    ACCENT_CYAN = RGBColor(0, 229, 255)  # Sci-Fi Cyan
    TEXT_LIGHT = RGBColor(240, 243, 246) # Crisp Off-White
    TEXT_MUTED = RGBColor(160, 174, 192) # Soft Gray
    STUB_BORDER = RGBColor(255, 99, 71)  # Tomato Accent for Stubs
    STUB_BG = RGBColor(45, 25, 30)       # Dark Red tint for Stub background

    img_dir = os.path.join(os.path.dirname(__file__), "..", "docs", "images")
    img_dir = os.path.abspath(img_dir)

    def set_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="PROJECT G1 OVERVIEW"):
        # Top bar accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_AMBER
        line.line.fill.background()

        # Category / Kicker text
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.733), Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.name = "Consolas"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN

        # Main Title text
        tx_box_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.733), Inches(0.7))
        tf_title = tx_box_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Trebuchet MS"
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_LIGHT

    def add_card(slide, left, top, width, height, title, body_bullets, bg_color=BG_CARD, border_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()

        tx_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3))
        tf = tx_box.text_frame
        tf.word_wrap = True

        if title:
            p_title = tf.paragraphs[0]
            p_title.text = title
            p_title.font.name = "Trebuchet MS"
            p_title.font.size = Pt(16)
            p_title.font.bold = True
            p_title.font.color.rgb = ACCENT_AMBER
            p_title.space_after = Pt(8)

        is_first = not bool(title)
        for bullet in body_bullets:
            p = tf.add_paragraph() if not is_first else tf.paragraphs[0]
            is_first = False
            p.text = "• " + bullet if isinstance(bullet, str) else "• " + bullet[0]
            p.font.name = "Calibri"
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_LIGHT
            p.space_after = Pt(4)
            if isinstance(bullet, tuple) and len(bullet) > 1:
                p_sub = tf.add_paragraph()
                p_sub.text = "   " + bullet[1]
                p_sub.font.name = "Calibri"
                p_sub.font.size = Pt(11)
                p_sub.font.color.rgb = TEXT_MUTED
                p_sub.space_after = Pt(6)

        return card

    def add_stub_box(slide, left, top, width, height, stub_title, stub_desc):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = STUB_BG
        box.line.color.rgb = STUB_BORDER
        box.line.width = Pt(2)

        tx_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), height - Inches(0.2))
        tf = tx_box.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = f"[STUB / PLACEHOLDER]: {stub_title}"
        p1.font.name = "Consolas"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = STUB_BORDER
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = stub_desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.italic = True
        p2.font.color.rgb = TEXT_MUTED
        p2.alignment = PP_ALIGN.CENTER
        return box

    def add_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1)

    # Accent banner box
    banner = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.2))
    banner.fill.solid()
    banner.fill.fore_color.rgb = BG_CARD
    banner.line.color.rgb = ACCENT_AMBER
    banner.line.width = Pt(2)

    tx_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(10.9), Inches(3.6))
    tf = tx_box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "RETRO FPS REVIVAL • 1998 GOLDSRC SPIRIT"
    p0.font.name = "Consolas"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_CYAN
    p0.space_after = Pt(14)

    p1 = tf.add_paragraph()
    p1.text = "PROJECT G1"
    p1.font.name = "Trebuchet MS"
    p1.font.size = Pt(46)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_AMBER
    p1.space_after = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = "Original Sci-Fi Story & Fully Scripted Procedural Asset Pipeline in Unity"
    p2.font.name = "Calibri"
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "Presenter: [STUB: Speaker Name / Team Lead]  |  Date: July 2026  |  Build Version: 2022.3 LTS"
    p3.font.name = "Consolas"
    p3.font.size = Pt(12)
    p3.font.color.rgb = TEXT_MUTED

    add_notes(slide1, "Welcome everyone to the Project G1 presentation. Introduce yourself and set the stage: Project G1 is a homage to classic 1998 retro action shooters, built with a modern programmatic pipeline in Unity and Blender.")

    # ==========================================
    # SLIDE 2: Executive Summary & Vision
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2)
    add_header(slide2, "Executive Summary & Core Pillars", "VISION & DESIGN GOALS")

    add_card(slide2, Inches(0.8), Inches(1.6), Inches(3.64), Inches(5.2),
             "1. Retro Aesthetic",
             [
                 ("1998 GoldSrc Heritage", "Captures the atmosphere of Half-Life 1 with low-poly geometry and industrial environments."),
                 ("Wear & Grime Pipeline", "Dynamic procedural texture masks multiplying color tints directly in standard shaders."),
                 ("Share Tech Mono HUD", "Authentic retro amber UI elements and tactile visual feedback.")
             ])

    add_card(slide2, Inches(4.84), Inches(1.6), Inches(3.64), Inches(5.2),
             "2. Authentic Gameplay",
             [
                 ("GoldSrc Physics Engine", "True air acceleration, 30-ups caps, bunnyhopping, coyote time, and HEV sprint."),
                 ("Tactical Combat & ADS", "5 core weapons + grenades with ADS zoom, recoil scaling, and secondary fire modes."),
                 ("Vehicle Integration", "24 arcade trucks across an 800x800m open sprawl for seamless traversal.")
             ])

    add_card(slide2, Inches(8.88), Inches(1.6), Inches(3.64), Inches(5.2),
             "3. 100% Procedural Engine",
             [
                 ("Blender Scripted Assets", "Zero manual 3D modeling — every character, weapon, vehicle, and building is code."),
                 ("Procedural Audio & Speech", "Math-synthesized retro sound effects + eSpeak NG automated dialogue pipeline."),
                 ("Self-Correcting Layouts", "Automated placement validation and dynamic interior room audio acoustics.")
             ])

    add_notes(slide2, "Highlight the three pillars: Retro Aesthetic (not cheap, but curated low-poly asset construction), Authentic Gameplay (Quake/HL movement), and 100% Procedural Engineering (code-driven pipeline).")

    # ==========================================
    # SLIDE 3: Narrative & World Bible
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3)
    add_header(slide3, "Narrative Loop & Character Lore", "STORY & WORLD BUILDING")

    add_card(slide3, Inches(0.8), Inches(1.6), Inches(6.5), Inches(5.2),
             "The Anomaly Loop Premise",
             [
                 ("The Threshold Event", "You play as a Senior Test Engineer surviving a catastrophic portal experiment failure at Corvus Facility."),
                 ("The Anchor & The Loop", "The Threshold isn't a door to a place—it's a loop of every past failure. You survive because you are the Anchor."),
                 ("The Aliens as Echoes", "Hostile creatures are the time-folded remains of past iterations trapped in the loop."),
                 ("The Auditor's Motive", "The man in the suit audits the loop for an interdimensional bureaucracy harvesting endless catastrophes."),
                 ("Branching Finale Choice", "Step through the Threshold ring (Loop resets) OR destroy the resonance emitters (Collapse loop & unmake the Anchor).")
             ])

    # Show Protagonist & Villain images
    protag_img = os.path.join(img_dir, "protagonist_front.png")
    villain_img = os.path.join(img_dir, "villain_front.png")

    if os.path.exists(protag_img):
        slide3.shapes.add_picture(protag_img, Inches(7.5), Inches(1.6), width=Inches(2.5))
    if os.path.exists(villain_img):
        slide3.shapes.add_picture(villain_img, Inches(10.1), Inches(1.6), width=Inches(2.4))

    # Caption box under images
    cap_box = slide3.shapes.add_textbox(Inches(7.5), Inches(6.3), Inches(5.0), Inches(0.6))
    tf_cap = cap_box.text_frame
    p_cap = tf_cap.paragraphs[0]
    p_cap.text = "Left: The Hazard Suit (Anchor)  |  Right: The Auditor"
    p_cap.font.name = "Consolas"
    p_cap.font.size = Pt(11)
    p_cap.font.color.rgb = ACCENT_AMBER
    p_cap.alignment = PP_ALIGN.CENTER

    add_notes(slide3, "Explain the core lore twist: The protagonist is the Anchor keeping the disaster alive, and the Auditor harvests the catastrophe. The player's first weapon (crowbar) is what destroys the loop at the end.")

    # ==========================================
    # SLIDE 4: Procedural Asset Pipeline (Blender to Unity)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4)
    add_header(slide4, "Code-Driven Asset & Modeling Pipeline", "PROCEDURAL ART PIPELINE")

    add_card(slide4, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2),
             "Blender Python Scripting Engine",
             [
                 ("Zero Manual Sculpting", "Characters built via elliptical tapered lofts (oval()), layered armor plates, and sealed helmets."),
                 ("Asymmetric Equipment", "Wrist computers, bracers, and uneven pouches avoid the artificial look of mirrored low-poly models."),
                 ("Wear & Grime Mask Baking", "Bakes dirt masks off ambient occlusion and height rather than pointiness, preventing edge bleaching."),
                 ("Standard Shader Integration", "White-to-grime mask multiplies cleanly over faction tint colors in Unity standard shaders."),
                 ("Rigging & Animation", "Automated bone weights and skeleton generation preserving classic 1998 animation clips.")
             ])

    add_stub_box(slide4, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
                 "BLENDER DEMO / PIPELINE DIAGRAM",
                 "Insert a video clip or flow diagram showing:\n\n1. python build_character.py execution\n2. Automated mesh generation in Blender\n3. Dirt mask texture baking\n4. Direct Unity FBX + Material import flow")

    add_notes(slide4, "Highlight why code-driven modeling matters: reproducible builds, instant variation tuning, automated UV mask baking without mudding albedo tints.")

    # ==========================================
    # SLIDE 5: Gameplay Mechanics & Movement Physics
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5)
    add_header(slide5, "Authentic GoldSrc Movement & Sprint Dynamics", "GAMEPLAY & PHYSICS")

    add_card(slide5, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2),
             "HL1 Physics & Sprint Metering",
             [
                 ("Quake-Lineage Acceleration", "Converts exact GoldSrc units to meters: strafe air steering, bunnyhop speed preservation."),
                 ("HEV Auxiliary Power Sprint", "Hold Shift to sprint at 12.6 m/s (up from 8.1 m/s), draining an 11-second auxiliary cell."),
                 ("Tactical Sprint Resource", "Meter locks out when depleted; UI bar pulses red. Metered movement turns open-field crossing into a tactical risk."),
                 ("ADS Zoom & Sway Tuning", "Hold RMB narrows FOV to 46°, reduces spread by 66%, dampens sway, and scales walk speed."),
                 ("Full Crouch & Coyote Time", "Smooth crouch height transitions and jump window forgiveness for precise platforming.")
             ])

    pov_img = os.path.join(img_dir, "unity_pov.png")
    if os.path.exists(pov_img):
        slide5.shapes.add_picture(pov_img, Inches(6.8), Inches(1.6), width=Inches(5.7))
        cap_box2 = slide5.shapes.add_textbox(Inches(6.8), Inches(6.2), Inches(5.7), Inches(0.6))
        tf_cap2 = cap_box2.text_frame
        p_cap2 = tf_cap2.paragraphs[0]
        p_cap2.text = "In-Engine Gameplay POV: Retro Amber HUD, Sprint Bar & Reticle"
        p_cap2.font.name = "Consolas"
        p_cap2.font.size = Pt(11)
        p_cap2.font.color.rgb = ACCENT_AMBER
        p_cap2.alignment = PP_ALIGN.CENTER
    else:
        add_stub_box(slide5, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
                     "GAMEPLAY SCREENSHOT / VIDEO",
                     "Insert HUD gameplay video showcasing:\n- Strafe jumping / Bunnyhopping\n- HEV auxiliary sprint depletion\n- Aim-Down-Sights (ADS) transition")

    add_notes(slide5, "Explain how movement feeling is identical to 1998 GoldSrc engine. Sprinting is metered so crossing open terrain requires strategic decision making.")

    # ==========================================
    # SLIDE 6: Weapon Arsenal & Equipment
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6)
    add_header(slide6, "Retro Firearms Arsenal & Secondary Fires", "WEAPONS & COMBAT")

    # Display Crowbar, Pistol, SMG images horizontally
    w_width = Inches(3.64)
    w_gap = Inches(0.4)

    weapons_data = [
        ("Crowbar & 9mm Pistol", ["Crowbar: Charged heavy swing (2.5x dmg + knockback).", "Pistol: 17-round mag, 3-round burst secondary fire."], "crowbar.png"),
        ("Shotgun & SMG 40mm", ["Shotgun: Per-shell reload, double-barrel blast.", "SMG: High rate of fire, integrated 40mm grenade launcher."], "smg.png"),
        (".357 Magnum & Grenades", ["Magnum: Heavy revolver cylinder state machine.", "Grenades: Cookable fuses, shockwave rings & shrapnel."], "pistol.png")
    ]

    for i, (title, bullets, img_filename) in enumerate(weapons_data):
        left_pos = Inches(0.8) + i * (w_width + w_gap)
        img_path = os.path.join(img_dir, img_filename)
        
        # Upper card with image
        if os.path.exists(img_path):
            slide6.shapes.add_picture(img_path, left_pos, Inches(1.6), width=w_width)

        # Lower card with specs
        add_card(slide6, left_pos, Inches(3.8), w_width, Inches(3.0), title, bullets)

    add_notes(slide6, "Detail the 5 weapons + grenades. Every weapon has animated slides/cylinders and custom secondary fire modes. Crowbar charged swing deals knockback.")

    # ==========================================
    # SLIDE 7: World Design & Vehicle Systems
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7)
    add_header(slide7, "800x800m Corvus Sprawl & Drivable Fleet", "WORLD & VEHICLES")

    add_card(slide7, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2),
             "Corvus Sprawl Base Architecture",
             [
                 ("50+ Enterable Interiors", "Hollow building structures with door cutouts, multi-floor storeys, rooftop catwalks."),
                 ("Military Base Layout", "280m runway, 3-storey control tower, 6 ammo igloos, tank workshop, trench combat zones."),
                 ("24 Parked Drivable Trucks", "Trucks spaced every 80m across ring roads. Approaching within 5m displays [E] DRIVE prompt."),
                 ("Arcade Vehicle Physics", "Predictable raycast + body-force physics optimized for high-speed ramp and berm terrain traversal."),
                 ("Interior Audio Acoustics", "G1InteriorSpace dynamically scales reverb duration (0.5s to 2.5s) and outdoor sound occlusion.")
             ])

    ov_img = os.path.join(img_dir, "unity_overview.png")
    if os.path.exists(ov_img):
        slide7.shapes.add_picture(ov_img, Inches(6.8), Inches(1.6), width=Inches(5.7))
        cap_box3 = slide7.shapes.add_textbox(Inches(6.8), Inches(6.2), Inches(5.7), Inches(0.6))
        tf_cap3 = cap_box3.text_frame
        p_cap3 = tf_cap3.paragraphs[0]
        p_cap3.text = "Overview map rendering: Corvus Sprawl district layout & interior spaces"
        p_cap3.font.name = "Consolas"
        p_cap3.font.size = Pt(11)
        p_cap3.font.color.rgb = ACCENT_AMBER
        p_cap3.alignment = PP_ALIGN.CENTER
    else:
        add_stub_box(slide7, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
                     "VEHICLE & MAP FOOTAGE",
                     "Insert gameplay video showing:\n- Driving a truck across Corvus Sprawl\n- Entering a workshop interior with acoustic reverb shift\n- Bio-scanner scanning nearby survivors")

    add_notes(slide7, "Explain the vehicle traversal system: 24 trucks ensure you are never more than 80m from a vehicle. Driving is noisy and draws enemy fire, balancing mobility vs stealth.")

    # ==========================================
    # SLIDE 8: Procedural Audio & Automated Voice Pipeline
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8)
    add_header(slide8, "Zero External Audio: Pure Math SFX & eSpeak NG", "PROCEDURAL AUDIO & VOICE")

    add_card(slide8, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2),
             "Math SFX & Automated Voice Casting",
             [
                 ("100% Procedural SFX", "Gunshots, explosions, pickup hums, door hydraulics, and footstep impacts generated from raw math."),
                 ("eSpeak NG Automated Voice", "Python script parses C# dialogue strings, generating voice clips per character profile."),
                 ("Script-As-Single-Source", "C# code defines dialogue lines. Rewording a line automatically mints a new clip with zero voice drift."),
                 ("Typewriter Pacing Sync", "Dialogue typewriter automatically scales display speed to match audio clip duration."),
                 ("Formant Fallback Engine", "Unvoiced fallback blips 6 formant-synthesized vowels, sounding like rapid radio chatter.")
             ])

    add_card(slide8, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
             "10 Distinct Voice Profiles",
             [
                 ("The Commander / Chief", "Low, barking, authoritative military cadence."),
                 ("Signals Tech", "Rapid-fire delivery at 3.4 words per second."),
                 ("The Auditor", "Unbothered, flat, sinister bureaucracy drawl."),
                 ("The Echo", "Croaking, slowed to 1.6 words/sec, dragged out of a time-folded body."),
                 ("Suit V.I.", "Monotone, even, synthetic tactical guidance.")
             ])

    add_notes(slide8, "Demonstrate the procedural speech: C# is the single source of truth for all text and voice audio clips. eSpeak NG generates voice files upon build.")

    # ==========================================
    # SLIDE 9: AI Tactical Engine & Pacing Director
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_bg(slide9)
    add_header(slide9, "GOAP Tactical AI & L4D2-Style Horde Pacing", "ENEMY AI & THREAT DIRECTOR")

    add_card(slide9, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2),
             "Squad Tactical AI (HECU Soldiers)",
             [
                 ("GOAP-Lite Decision Planner", "Fighters evaluate threat distance, cover safety, ammo, and squad flanking maneuvers."),
                 ("Geometry-Derived Firing Cover", "Map generator extracts fire steps, sandbag walls, and slit openings directly from geometry."),
                 ("Dynamic Cover Re-evaluation", "Soldiers claim cover, fire from slits, and abandon compromised positions when flanked."),
                 ("Pruned Reachability Check", "Post-bake validation removes buried/unreachable nodes so AI never gets stuck.")
             ])

    add_card(slide9, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
             "Horde Pacing & Boss Mechanics",
             [
                 ("L4D2-Style ThreatDirector", "Paces zombie & alien spawns dynamically based on player stress and location."),
                 ("Bio-Scanner Radar (Q)", "Sweeps 150m radius, marking survivor contacts and reporting bearings to unknown signals."),
                 ("HECU Gunship Boss (Level 2)", "Strafing machine-gun sweeps, 3-rocket salvos, and destructible rotor physics."),
                 ("Threshold Anomaly Boss (Level 3)", "Multi-phase portal breach fight requiring resonance emitter destruction.")
             ])

    add_notes(slide9, "Explain AI intelligence: HECU soldiers use real geometry cover (sandbag slits, trench steps). ThreatDirector prevents monotonous enemy spawning.")

    # ==========================================
    # SLIDE 10: Campaign Structure & 3-Level Flow
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_bg(slide10)
    add_header(slide10, "3-Level Campaign Arc & Chapter Spine", "CAMPAIGN & LEVEL DESIGN")

    l_width = Inches(3.64)
    l_gap = Inches(0.4)

    add_card(slide10, Inches(0.8), Inches(1.6), l_width, Inches(5.2),
             "Level 1: Corvus Annex",
             [
                 ("Facility Breach", "Locker Room → Lab Corridor → Control Room → Industrial Hall."),
                 ("HECU Ambush", "First combat encounter against military containment squads."),
                 ("Alien Breach Zone", "Xen portal leaks and horror atmosphere leading to the emergency elevator.")
             ])

    add_card(slide10, Inches(0.8) + l_width + l_gap, Inches(1.6), l_width, Inches(5.2),
             "Level 2: Quarantine Zone",
             [
                 ("Outdoor Complex", "Toxic hazard zones, jump pads, and long-range squad engagements."),
                 ("Vehicle Warfare", "Drive trucks across 800m open ground between districts."),
                 ("Gunship Boss Fight", "Destructible rotor battle against heavy assault helicopter.")
             ])

    add_card(slide10, Inches(0.8) + 2 * (l_width + l_gap), Inches(1.6), l_width, Inches(5.2),
             "Level 3: Threshold Anomaly",
             [
                 ("Portal Breach", "Dimensional collapse chamber filled with time-folded echoes."),
                 ("Resonance Emitters", "Destroy three emitter towers holding the anomaly ring open."),
                 ("The Final Choice", "Step through to reset the loop OR destroy emitters to unmake the Anchor.")
             ])

    add_notes(slide10, "Walk through the campaign progression: from indoor horror facility (L1) to open vehicle combat (L2) to dimensional anomaly climax (L3).")

    # ==========================================
    # SLIDE 11: Developer & QA Tooling Sandbox
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_bg(slide11)
    add_header(slide11, "Built-in Developer Sandbox & Telemetry", "QA & OBSERVABILITY")

    add_card(slide11, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2),
             "Developer Controls & Testing Range",
             [
                 ("Weapon Testing Range", "Dedicated scene for testing weapon spread, recoil, and damage values."),
                 ("God Mode & Infinite Ammo (G)", "Unkillable state tied with auto-replenishing magazines and HUD infinity symbol."),
                 ("3D Fly / Noclip Mode (V)", "Free camera traversal across the 800m map to inspect geometry and AI paths."),
                 ("Mob Spawner Toolbox (TAB)", "Dynamically spawn Zombies, HECU Soldiers, Squads, Hordes, or Bosses on demand."),
                 ("Instant Reset (F5)", "Reload active scene instantly for rapid playtesting iteration.")
             ])

    add_card(slide11, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
             "Observability & Build Verification",
             [
                 ("F3 Telemetry Overlay", "Live FPS, draw calls, memory usage, AI state counters, and ThreatDirector stress levels."),
                 ("AI State Gizmos", "In-editor visual rays for target line-of-sight, cover node claims, and pathfinding."),
                 ("G1VerifyBuild Headless Tool", "Automated headless check verifying map contact coordinates aren't inside solid blocks."),
                 ("G1SelfTest Play Mode Tool", "Automated arm playtest script testing map relocation and door triggers.")
             ])

    add_notes(slide11, "Emphasize QA readiness: The game includes built-in observability tools (F3 overlay, TAB mob spawner, headless automated verification) to ensure build quality.")

    # ==========================================
    # SLIDE 12: Technical Architecture & System Requirements
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_bg(slide12)
    add_header(slide12, "Unity Engine Integration & Dependencies", "TECHNICAL STACK")

    add_card(slide12, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2),
             "Engine & Build Baseline",
             [
                 ("Unity 2022.3.62f3 LTS", "Built-in Render Pipeline for classic 1998 lighting and retro shader compatibility."),
                 ("Classic Input Manager", "Standardized keyboard/mouse binding mapping with smooth mouse look smoothing."),
                 ("Git LFS Tracking", "All model binaries (*.fbx) tracked via Git LFS to prevent empty pointer imports."),
                 ("Procedural Scene Builders", "G1SceneBuilder, G1CampaignBuilders, G1MenuBuilder rebuild scenes from scratch code.")
             ])

    add_card(slide12, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
             "External Pipeline Tools",
             [
                 ("Blender 4.x / 5.x", "Used in headless mode to run python scripts for characters, weapons, vehicles, and maps."),
                 ("eSpeak NG Speech Engine", "Synthesizes character dialogue audio directly into committed .wav files."),
                 ("Cross-Session Save System", "JSON-serialized player progress stored in persistentDataPath.")
             ])

    add_notes(slide12, "Review tech stack requirements: Unity 2022.3 LTS with Built-in Render Pipeline, Blender 4.x/5.x for asset generation, eSpeak NG for voice audio.")

    # ==========================================
    # SLIDE 13: Project Roadmap & Achievements
    # ==========================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_bg(slide13)
    add_header(slide13, "Current Status & Future Milestones", "ROADMAP & MILESTONES")

    add_card(slide13, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2),
             "Completed Features [COMPLETED]",
             [
                 ("✓ Full 5-Weapon Arsenal & ADS", "Crowbar, Pistol, Shotgun, SMG, Magnum + grenades with secondary fire."),
                 ("✓ 3-Level Campaign Generation", "Corvus Annex → Quarantine Zone → Threshold Boss Arena."),
                 ("✓ Drivable Truck Fleet & Sprawl", "24 arcade trucks across 800m map with enterable interiors."),
                 ("✓ Procedural Audio & Voice Clips", "Math SFX + eSpeak NG synthesized voice actor pipeline."),
                 ("✓ Branching Cinematic Finale", "Stabilize or Collapse the loop ending choices.")
             ])

    add_card(slide13, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
             "Future Milestones [IN PROGRESS]",
             [
                 ("▶ Advanced Boss Mechanics", "Multi-phase alien boss attack patterns and shield disruption mechanics."),
                 ("▶ Serialized State Saving", "Complete cross-session inventory and objective save serialization."),
                 ("▶ Procedural Seed Modding", "Exporting/importing custom level seeds for infinite map generation."),
                 ("▶ Community Modding Kit", "Exposing Blender python asset builder parameters to mod creators.")
             ])

    add_notes(slide13, "Summarize roadmap accomplishments: Core campaign, weapons, trucks, audio, and ending choices are done. Future focus is on boss mechanics and seed modding.")

    # ==========================================
    # SLIDE 14: Presentation Stubs & Live Demo Checklist
    # ==========================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_bg(slide14)
    add_header(slide14, "Guidelines for Presenter & Live Demonstration", "LIVE DEMO & STUBS GUIDE")

    add_stub_box(slide14, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2),
                 "LIVE DEMO STEPS (3 MINS)",
                 "1. Launch Unity 2022.3 LTS & open TestScene\n2. Press Play and showcase HL1 movement (strafe jump / sprint)\n3. Press [TAB] to spawn HECU Squad vs Zombies\n4. Enter truck with [E] and drive through workshop\n5. Demonstrate ADS RMB aim and secondary fire burst")

    add_stub_box(slide14, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
                 "SLIDE CUSTOMIZATION STUBS",
                 "Replace the following placeholders before client/team presentation:\n\n• Slide 1: Presenter Name & Team Title\n• Slide 4: Optional video of Blender script generating character\n• Slide 11: Real benchmark performance graph (FPS on target GPU)\n• Q&A Slide: Contact info / Git repository link")

    add_notes(slide14, "Use this slide as a reference during presentation rehearsal. Conduct the 3-minute live demo showing movement, spawner, and truck driving.")

    # ==========================================
    # SLIDE 15: Conclusion & Q&A
    # ==========================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_bg(slide15)

    card_end = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.733), Inches(5.2))
    card_end.fill.solid()
    card_end.fill.fore_color.rgb = BG_CARD
    card_end.line.color.rgb = ACCENT_AMBER
    card_end.line.width = Pt(2)

    tx_box_end = slide15.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(10.9), Inches(4.4))
    tf_end = tx_box_end.text_frame
    tf_end.word_wrap = True

    pe0 = tf_end.paragraphs[0]
    pe0.text = "PROJECT G1 • RETRO FPS"
    pe0.font.name = "Consolas"
    pe0.font.size = Pt(14)
    pe0.font.color.rgb = ACCENT_CYAN
    pe0.space_after = Pt(14)

    pe1 = tf_end.add_paragraph()
    pe1.text = "THANK YOU! QUESTIONS & DEMO"
    pe1.font.name = "Trebuchet MS"
    pe1.font.size = Pt(40)
    pe1.font.bold = True
    pe1.font.color.rgb = ACCENT_AMBER
    pe1.space_after = Pt(16)

    pe2 = tf_end.add_paragraph()
    pe2.text = "Repository: https://github.com/Thundercok/Project-G1"
    pe2.font.name = "Consolas"
    pe2.font.size = Pt(16)
    pe2.font.color.rgb = TEXT_LIGHT
    pe2.space_after = Pt(8)

    pe3 = tf_end.add_paragraph()
    pe3.text = "Documentation: docs/architecture.md  |  docs/story.md  |  docs/asset-pipeline.md"
    pe3.font.name = "Consolas"
    pe3.font.size = Pt(14)
    pe3.font.color.rgb = TEXT_MUTED
    pe3.space_after = Pt(24)

    pe4 = tf_end.add_paragraph()
    pe4.text = "[STUB: Insert Team Contact Info / QR Code / Live Demo Links Here]"
    pe4.font.name = "Calibri"
    pe4.font.size = Pt(14)
    pe4.font.italic = True
    pe4.font.color.rgb = STUB_BORDER

    add_notes(slide15, "Thank the audience and invite questions or live playtesting of Project G1!")

    output_path = os.path.join(os.path.dirname(__file__), "..", "Project_G1_Presentation.pptx")
    output_path = os.path.abspath(output_path)
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_deck()
